"""
Per-format text extraction. Routes to the best extractor for each file type.
Falls back to unstructured for anything exotic.
"""
import io
from typing import Optional

from app.rag.detector import InputMeta


async def extract_text(file_bytes: bytes, meta: InputMeta, filename: str) -> str:
    """Route file to correct extractor. Returns raw extracted text."""
    ft = meta.file_type

    if ft == "pdf":
        return _extract_pdf(file_bytes, meta.is_scanned_pdf)
    if ft == "docx":
        return _extract_docx(file_bytes)
    if ft in ("doc",):
        return _extract_unstructured(file_bytes, filename)
    if ft == "xlsx":
        return _extract_xlsx(file_bytes)
    if ft == "csv":
        return _extract_csv(file_bytes)
    if ft == "txt":
        return file_bytes.decode("utf-8", errors="replace")
    if ft == "html":
        return _extract_html(file_bytes)
    if ft == "image":
        return _extract_image_ocr(file_bytes)
    if ft == "pptx":
        return _extract_pptx(file_bytes)

    # Unknown — try unstructured as catch-all
    return _extract_unstructured(file_bytes, filename)


def _extract_pdf(file_bytes: bytes, is_scanned: bool) -> str:
    if is_scanned:
        return _pdf_ocr(file_bytes)
    return _pdf_text(file_bytes)


def _pdf_text(file_bytes: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
            # Extract tables as text too
            for table in page.extract_tables():
                rows = ["\t".join(str(c or "") for c in row) for row in table]
                text_parts.append("\n".join(rows))
    return "\n\n".join(text_parts)


def _pdf_ocr(file_bytes: bytes) -> str:
    """
    OCR a scanned PDF.

    Needs the tesseract and poppler binaries, not just the Python wrappers, so
    it cannot be fixed by adding a requirement. The message is written to be
    read by a person: the job error string is shown in the dashboard, and
    "No module named pdf2image" tells an estate agent nothing about what to do.
    """
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError as exc:
        raise RuntimeError(
            "This PDF has no selectable text, so it needs OCR - which is not "
            "enabled on this deployment. Upload a text-based PDF, or paste the "
            "content in directly."
        ) from exc
    images = convert_from_bytes(file_bytes, dpi=200)
    pages = [pytesseract.image_to_string(img) for img in images]
    return "\n\n".join(pages)


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            parts.append(row_text)
    return "\n".join(parts)


def _extract_xlsx(file_bytes: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(v) if v is not None else "" for v in row)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_csv(file_bytes: bytes) -> str:
    """
    CSV to tab-separated text.

    Deliberately the standard library rather than pandas: this is the only place
    pandas was used, for one read-and-stringify, and it carries ~50 MB of
    dependency into every deploy for the privilege.
    """
    import csv
    text = file_bytes.decode("utf-8", errors="replace")
    rows = list(csv.reader(io.StringIO(text)))
    return "\n".join("\t".join(cell.strip() for cell in row) for row in rows if any(row))


def _extract_html(file_bytes: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _extract_image_ocr(file_bytes: bytes) -> str:
    """Read text out of an image. Same binary requirement as _pdf_ocr."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError(
            "Reading text out of images needs OCR, which is not enabled on this "
            "deployment. Paste the content directly instead."
        ) from exc
    img = Image.open(io.BytesIO(file_bytes))
    return pytesseract.image_to_string(img)


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"Slide {i}:\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _extract_unstructured(file_bytes: bytes, filename: str) -> str:
    try:
        from unstructured.partition.auto import partition
        elements = partition(file=io.BytesIO(file_bytes), metadata_filename=filename)
        return "\n\n".join(str(el) for el in elements if str(el).strip())
    except ImportError:
        return (
            "[This file type is not supported on this deployment. "
            "Supported: PDF, DOCX, XLSX, CSV, HTML, PPTX and plain text.]"
        )
    except Exception as e:
        return f"[Extraction failed: {e}]"
